#!/usr/bin/env python3
"""
Very small MD -> PPTX converter for onboarding decks.

Rules (intentionally simple):
- Slide 1: Title slide from first H1 (# ...)
- Each H2 (## ...) starts a new slide
- H3 (### ...) becomes a bold-ish line (prefix) within the slide
- Bullet lines (- / * ) become bullet paragraphs
- Quotes/blank lines are ignored
"""

from __future__ import annotations

import argparse
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt


H1_RE = re.compile(r"^\s*#\s+(?P<t>.+?)\s*$")
H2_RE = re.compile(r"^\s*##\s+(?P<t>.+?)\s*$")
H3_RE = re.compile(r"^\s*###\s+(?P<t>.+?)\s*$")
BULLET_RE = re.compile(r"^\s*[-*]\s+(?P<t>.+?)\s*$")


def _clean_inline(text: str) -> str:
    # remove inline markdown that commonly appears in this doc
    text = re.sub(r"`([^`]+)`", r"\1", text)
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    text = re.sub(r"\*([^*]+)\*", r"\1", text)
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)  # links -> anchor text
    text = text.replace("<br>", " ")
    return text.strip()


def md_to_pptx(md_path: Path, pptx_path: Path) -> None:
    md = md_path.read_text(encoding="utf-8").splitlines()

    prs = Presentation()

    title = None
    for line in md:
        m = H1_RE.match(line)
        if m:
            title = _clean_inline(m.group("t"))
            break
    if not title:
        title = md_path.stem

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if slide.placeholders and len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Auto-generated from markdown"

    current_slide = None
    body_tf = None
    last_was_text = False

    def new_section(section_title: str):
        nonlocal current_slide, body_tf, last_was_text
        current_slide = prs.slides.add_slide(prs.slide_layouts[1])
        current_slide.shapes.title.text = _clean_inline(section_title)
        body_tf = current_slide.shapes.placeholders[1].text_frame
        body_tf.clear()
        last_was_text = False

    def add_para(text: str, level: int = 0, font_size: int = 20, bold: bool = False):
        nonlocal body_tf, last_was_text
        if body_tf is None:
            return
        p = body_tf.paragraphs[0] if (not body_tf.paragraphs or not last_was_text) else body_tf.add_paragraph()
        p.text = _clean_inline(text)
        p.level = level
        for run in p.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
        last_was_text = True

    for raw in md:
        line = raw.strip()
        if not line:
            continue
        if line.startswith(">"):
            continue
        if line.startswith("---"):
            continue

        m2 = H2_RE.match(raw)
        if m2:
            new_section(m2.group("t"))
            continue

        m3 = H3_RE.match(raw)
        if m3 and body_tf is not None:
            add_para(m3.group("t"), level=0, font_size=22, bold=True)
            continue

        mb = BULLET_RE.match(raw)
        if mb and body_tf is not None:
            add_para(mb.group("t"), level=1, font_size=18, bold=False)
            continue

        # Skip code fences, images, callouts, and html-like blocks
        if raw.startswith("```") or raw.startswith("<"):
            continue

        # normal text line
        if body_tf is None:
            continue
        add_para(raw, level=0, font_size=18, bold=False)

    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(pptx_path))


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--input", required=True, help="input markdown path")
    ap.add_argument("--output", required=True, help="output pptx path")
    args = ap.parse_args()

    md_to_pptx(Path(args.input), Path(args.output))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())


